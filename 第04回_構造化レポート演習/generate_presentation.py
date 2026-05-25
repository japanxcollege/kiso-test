#!/usr/bin/env python3
"""ISM/DEMATEL レポート用 PowerPoint 生成"""
import json
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

BASE = Path(__file__).parent
RESULTS = json.loads((BASE / "results.json").read_text(encoding="utf-8"))
ITEMS = RESULTS["items"]
HIER = RESULTS["hierarchy"]
SK = RESULTS["skeleton"]
DpR = RESULTS["DpR"]
DmR = RESULTS["DmR"]


def add_title_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_bullet_slide(prs, title, bullets, font_size=18):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)


def add_table_slide(prs, title, headers, rows, col_widths=None):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    rows_n = len(rows) + 1
    cols_n = len(headers)
    left, top, width, height = Inches(0.5), Inches(1.4), Inches(9), Inches(0.35 * rows_n)
    table = slide.shapes.add_table(rows_n, cols_n, left, top, width, height).table
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(52, 152, 219)
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.font.size = Pt(11)
            p.font.bold = True
    for i, row in enumerate(rows, 1):
        for j, val in enumerate(row):
            cell = table.cell(i, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)


def matrix_rows(mat, labels):
    rows = []
    for i, row in enumerate(mat):
        rows.append([f"{i+1}:{labels[i][:8]}"] + [str(v) for v in row])
    return rows


def build():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "ISM法・DEMATEL法・因果ループ図による\n問題構造化レポート",
        "システム工学基礎（第2回課題）\n03-240974  高木俊輔\n2026年5月16日提出",
    )

    # 設問1
    add_bullet_slide(
        prs,
        "設問1：解決したい問題",
        [
            "【問題の明示】",
            "大学4年次において、必修・選択科目の単位取得とGPAを「安定的に」確保する。",
            "",
            "【現状と理想のギャップ】",
            "・理想：全科目でC以上、留年・単位不足リスクを排除したGPA 2.5以上",
            "・現状：バイト・課外と学習の競合、睡眠不足、計画性不足で成績がぶれやすい",
            "",
            "【境界設定（認識論的境界）】",
            "・含む：学習行動、生活リズム、授業内の理解、評価（試験・レポート）",
            "・含まない：就活面接スキル、大学制度変更、経済的支援（奨学金等）",
            "",
            "【選定理由】",
            "第1回小演習で挙げた「単位・睡眠・バイト」の資源競合を、構造分析で定式化する。",
        ],
        16,
    )

    # 設問2 - 要素
    add_bullet_slide(
        prs,
        "設問2：要素抽出（8要素）",
        [f"{i+1}. {name}" for i, name in enumerate(ITEMS)]
        + [
            "",
            "専門家（本人）による一対比較で「i が j に影響するか」を判定（0/1）。",
            "推移律を仮定：間接影響は直接判定を省略可能。",
        ],
        18,
    )

    # 隣接行列
    headers = ["要素"] + [str(i + 1) for i in range(len(ITEMS))]
    add_table_slide(prs, "設問2：隣接行列（構造行列）", headers, matrix_rows(RESULTS["A"], ITEMS))

    # メンタルモデル
    add_bullet_slide(
        prs,
        "設問2：メンタルモデルとの比較（反復）",
        [
            "【初期メンタルモデル】",
            "「学習時間↑→成績↑」「出席↑→理解↑」が中心。バイトは時間を奪う要因程度。",
            "",
            "【第1反復】可到達行列・階層化後",
            "・GPAは最上位、バイト・課外が最下位（根本原因側）と判明",
            "・睡眠→意欲→学習時間の連鎖が中間層に現れた → 睡眠の辺を追加",
            "",
            "【第2反復】骨格行列（スケルトン）抽出後",
            "・学習時間→試験質の直接辺は、理解を介する推移で冗長 → 削除",
            "・バイト→学習時間の直接辺も、睡眠・意欲経由で冗長 → 削除",
            "・残った7本の辺が「本質的な直接因果」のスケルトン",
            "",
            "→ メンタルモデルは「時間投入＝成績」と短絡していた点を修正。",
        ],
        15,
    )

    # 階層
    hier_lines = []
    for lv, elems in enumerate(HIER):
        names = ", ".join(f"{e}:{ITEMS[e-1]}" for e in elems)
        hier_lines.append(f"レベル{lv+1}（結果側↑）: {names}")
    add_bullet_slide(prs, "設問2：ISM階層構造", hier_lines + ["", "解釈：下位＝原因系、上位＝結果系。単位・GPAはシステムの「出口」。"], 14)

    # 骨格
    sk_edges = []
    for i in range(len(ITEMS)):
        for j in range(len(ITEMS)):
            if SK[i][j]:
                sk_edges.append(f"{i+1}→{j+1}: {ITEMS[i]} → {ITEMS[j]}")
    add_bullet_slide(prs, "設問2：骨格行列（スケルトン）", sk_edges, 14)

    # 可到達（要約）
    add_bullet_slide(
        prs,
        "設問2：可到達行列（抜粋・解釈）",
        [
            "全要素間で到達可能（連結）。ブール積の収束後、",
            "レベル7（バイト）からレベル1（GPA）へ一方向の因果連鎖が確認された。",
            "",
            "R∩A=R を満たす要素から順に抽出 → 7段階の階層。",
            "計算：2026_ISM-DEMATEL.xlsm / ism_dematel.html で検証。",
        ],
        16,
    )

    # 設問3 DEMATEL
    add_bullet_slide(
        prs,
        "設問3：DEMATEL（直接影響 0〜4）",
        [
            "ISMの0/1に加え、影響の「強さ」を専門家評価（0=なし, 4=非常に強い）。",
            "例：バイト→学習時間=4、睡眠→試験質=3、試験質→GPA=4",
            "正規化 N = X / max(行和, 列和)、総合影響 T = N(I−N)⁻¹",
        ],
        16,
    )

    headers_d = ["要素", "D+R(関連度)", "D-R(影響度)", "分類"]
    rows_d = []
    for i in range(len(ITEMS)):
        role = "原因系" if DmR[i] > 0.05 else ("結果系" if DmR[i] < -0.05 else "媒介")
        rows_d.append([f"{i+1}.{ITEMS[i]}", f"{DpR[i]:.3f}", f"{DmR[i]:.3f}", role])
    add_table_slide(prs, "設問3：DEMATEL 影響度・関連度", headers_d, rows_d)

    add_bullet_slide(
        prs,
        "設問3：影響関係の強さの議論",
        [
            "【原因系（D−R > 0）】",
            "・バイト・課外（+1.18）：他要素への純送出が最大。時間・睡眠を圧迫する起点。",
            "・睡眠・体調（+0.69）：意欲・理解への間接影響が大きい基盤要因。",
            "・学習環境（+0.59）：時間配分の効率を左右するが、単独では出口に届きにくい。",
            "",
            "【結果系（D−R < 0）】",
            "・試験・レポートの質（−1.39）：関連度(D+R)も最大→システムの「ハブ」",
            "・GPA・単位（−1.00）：純粋な結果変数。他からの影響を最も受ける。",
            "",
            "【ISMとの整合】",
            "階層の最下位＝DEMATEL原因系、最上位＝結果系で一致。",
            "媒介要素（学習時間・理解）は関連度が高く、改善のレバー候補。",
        ],
        14,
    )

    # 設問4 因果ループ
    add_bullet_slide(
        prs,
        "設問4：因果ループ図（CLD）",
        [
            "骨格構造を動的フィードバックとして再記述（第3回講義の記法）。",
            "",
            "【強化ループ R1：好循環】",
            "学習意欲↑─(S)→学習時間↑─(S)→理解↑─(S)→成績↑─(S)→自信↑─(S)→意欲↑",
            "",
            "【バランスループ B1：バイトの負荷】",
            "バイト時間↑─(O)→学習時間↓─(O)→成績↓─(S)→経済不安↑─(S)→バイト時間↑",
            "",
            "【バランスループ B2：睡眠不足】",
            "睡眠不足↑─(O)→集中↓─(O)→成績↓─(S)→焦り↑─(O)→睡眠時間↓",
            "",
            "【本質の把握】",
            "単位問題は「努力不足」だけでなく、B1・B2の負のループが",
            "R1を抑制している構造。介入はバイト量・睡眠の境界調整が有効。",
        ],
        13,
    )

    add_bullet_slide(
        prs,
        "まとめ",
        [
            "ISM：階層と骨格で「何が何に先行するか」を整理",
            "DEMATEL：原因系/結果系と影響の強さを定量比較",
            "CLD：時間経過でのフィードバック（好循環・悪循環）を可視化",
            "",
            "問題の本質：成績（GPA）以前に、生活リズムと資源配分の構造がボトルネック。",
            "",
            "使用ツール：2026_ISM-DEMATEL.xlsm（同等計算：ism_dematel.html）",
        ],
        18,
    )

    out = BASE / "20260516_構造化レポート.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
